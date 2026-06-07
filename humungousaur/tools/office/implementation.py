from __future__ import annotations

from datetime import datetime, timezone
from html import escape
import json
from pathlib import Path
import re
from typing import Any
from uuid import uuid4
from zipfile import ZIP_DEFLATED, ZipFile
import xml.etree.ElementTree as ET

from humungousaur.config import AgentConfig
from humungousaur.schemas import ActionStatus, RiskLevel, ToolResult
from humungousaur.tools.base import Tool, object_input_schema


OFFICE_MAX_PARAGRAPHS = 500
OFFICE_MAX_TABLE_ROWS = 500
OFFICE_MAX_SLIDES = 80
OFFICE_MAX_TEXT_CHARS = 20_000


class DocxDocumentCreateTool(Tool):
    def __init__(self) -> None:
        super().__init__(
            name="docx_document_create",
            description=(
                "Create a native DOCX artifact from structured title, headings, paragraphs, bullets, and tables. "
                "This is a Humungousaur-owned OOXML writer and does not import upstream document scripts."
            ),
            risk_level=RiskLevel.MEDIUM,
            input_schema=object_input_schema(
                {
                    "filename": {"type": "string", "description": "Output filename under data_dir/documents."},
                    "title": {"type": "string", "description": "Document title."},
                    "sections": {
                        "type": "array",
                        "items": {"type": "object"},
                        "description": "Section specs with heading, paragraphs, bullets, and optional tables.",
                    },
                    "reason": {"type": "string", "description": "Why this DOCX artifact should be created."},
                },
                required=["title", "sections", "reason"],
            ),
            capability_group="office",
        )

    def execute(self, tool_input: dict[str, Any], config: AgentConfig) -> ToolResult:
        normalized = config.normalized()
        title = " ".join(str(tool_input.get("title") or "").split())
        sections = tool_input.get("sections")
        if not title:
            return ToolResult(self.name, ActionStatus.FAILED, self.risk_level, "Document title is required.")
        if not isinstance(sections, list) or not sections:
            return ToolResult(self.name, ActionStatus.FAILED, self.risk_level, "At least one document section is required.")
        filename = _safe_filename(str(tool_input.get("filename") or "humungousaur-document.docx"), ".docx")
        path = (normalized.data_dir / "documents" / filename).resolve()
        if not _is_within(path, normalized.allowed_write_roots):
            return ToolResult(self.name, ActionStatus.BLOCKED, self.risk_level, "DOCX path is outside allowed write roots.")
        summary = _docx_summary(title, sections)
        if config.dry_run:
            return ToolResult(self.name, ActionStatus.SKIPPED, self.risk_level, f"Dry run: would create DOCX document {path}.", {"path": str(path), **summary})
        path.parent.mkdir(parents=True, exist_ok=True)
        _write_docx(path, title=title, sections=sections)
        return ToolResult(
            self.name,
            ActionStatus.SUCCEEDED,
            self.risk_level,
            f"Created DOCX document {path}.",
            {"path": str(path), **summary, "source": "docx_document_create"},
        )


class DocxDocumentInspectTool(Tool):
    def __init__(self) -> None:
        super().__init__(
            name="docx_document_inspect",
            description="Inspect a local DOCX artifact for paragraphs, tables, headings, and sample text.",
            risk_level=RiskLevel.LOW,
            input_schema=object_input_schema(
                {
                    "path": {"type": "string", "description": "Workspace-relative or allowed absolute DOCX path."},
                    "sample_paragraphs": {"type": "integer", "minimum": 1, "maximum": 50},
                },
                required=["path"],
            ),
            capability_group="office",
        )

    def execute(self, tool_input: dict[str, Any], config: AgentConfig) -> ToolResult:
        normalized = config.normalized()
        path = _resolve_office_path(normalized, str(tool_input.get("path") or ""), subdir="documents", suffix=".docx")
        if not _is_within(path, normalized.allowed_read_roots + normalized.allowed_write_roots):
            return ToolResult(self.name, ActionStatus.BLOCKED, self.risk_level, "DOCX path is outside allowed roots.")
        if not path.exists() or path.suffix.lower() != ".docx":
            return ToolResult(self.name, ActionStatus.FAILED, self.risk_level, "DOCX file does not exist.")
        try:
            payload = _inspect_docx(path, sample_paragraphs=max(1, min(int(tool_input.get("sample_paragraphs") or 10), 50)))
        except Exception as exc:
            return ToolResult(self.name, ActionStatus.FAILED, self.risk_level, "DOCX file could not be inspected.", error=str(exc))
        return ToolResult(
            self.name,
            ActionStatus.SUCCEEDED,
            self.risk_level,
            f"Inspected DOCX document {path}.",
            {"path": str(path), **payload, "source": "docx_document_inspect"},
        )


class PptxDeckCreateTool(Tool):
    def __init__(self) -> None:
        super().__init__(
            name="pptx_deck_create",
            description=(
                "Create a native PPTX deck artifact from structured slide specs with titles, bullets, and speaker notes text. "
                "This uses Humungousaur-owned schemas and local Python PowerPoint support when available."
            ),
            risk_level=RiskLevel.MEDIUM,
            input_schema=object_input_schema(
                {
                    "filename": {"type": "string", "description": "Output filename under data_dir/presentations."},
                    "title": {"type": "string", "description": "Deck title."},
                    "slides": {
                        "type": "array",
                        "items": {"type": "object"},
                        "description": "Slide specs with title, bullets, body, and notes.",
                    },
                    "reason": {"type": "string", "description": "Why this PPTX deck should be created."},
                },
                required=["title", "slides", "reason"],
            ),
            capability_group="office",
        )

    def execute(self, tool_input: dict[str, Any], config: AgentConfig) -> ToolResult:
        try:
            from pptx import Presentation
        except Exception as exc:
            return ToolResult(self.name, ActionStatus.FAILED, self.risk_level, "python-pptx is required for PPTX creation.", error=str(exc))
        normalized = config.normalized()
        title = " ".join(str(tool_input.get("title") or "").split())
        slides = tool_input.get("slides")
        if not title:
            return ToolResult(self.name, ActionStatus.FAILED, self.risk_level, "Deck title is required.")
        if not isinstance(slides, list) or not slides:
            return ToolResult(self.name, ActionStatus.FAILED, self.risk_level, "At least one slide spec is required.")
        if len(slides) > OFFICE_MAX_SLIDES:
            return ToolResult(self.name, ActionStatus.BLOCKED, self.risk_level, "PPTX slide count exceeds safety limit.")
        filename = _safe_filename(str(tool_input.get("filename") or "humungousaur-deck.pptx"), ".pptx")
        path = (normalized.data_dir / "presentations" / filename).resolve()
        if not _is_within(path, normalized.allowed_write_roots):
            return ToolResult(self.name, ActionStatus.BLOCKED, self.risk_level, "PPTX path is outside allowed write roots.")
        summary = _pptx_summary(title, slides)
        if config.dry_run:
            return ToolResult(self.name, ActionStatus.SKIPPED, self.risk_level, f"Dry run: would create PPTX deck {path}.", {"path": str(path), **summary})
        presentation = Presentation()
        _add_title_slide(presentation, title)
        for slide_spec in slides:
            if isinstance(slide_spec, dict):
                _add_content_slide(presentation, slide_spec)
        path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(path)
        return ToolResult(
            self.name,
            ActionStatus.SUCCEEDED,
            self.risk_level,
            f"Created PPTX deck {path}.",
            {"path": str(path), **summary, "source": "pptx_deck_create"},
        )


class PptxDeckInspectTool(Tool):
    def __init__(self) -> None:
        super().__init__(
            name="pptx_deck_inspect",
            description="Inspect a local PPTX deck artifact for slide count, titles, and sample text.",
            risk_level=RiskLevel.LOW,
            input_schema=object_input_schema(
                {
                    "path": {"type": "string", "description": "Workspace-relative or allowed absolute PPTX path."},
                    "sample_slides": {"type": "integer", "minimum": 1, "maximum": 50},
                },
                required=["path"],
            ),
            capability_group="office",
        )

    def execute(self, tool_input: dict[str, Any], config: AgentConfig) -> ToolResult:
        try:
            from pptx import Presentation
        except Exception as exc:
            return ToolResult(self.name, ActionStatus.FAILED, self.risk_level, "python-pptx is required for PPTX inspection.", error=str(exc))
        normalized = config.normalized()
        path = _resolve_office_path(normalized, str(tool_input.get("path") or ""), subdir="presentations", suffix=".pptx")
        if not _is_within(path, normalized.allowed_read_roots + normalized.allowed_write_roots):
            return ToolResult(self.name, ActionStatus.BLOCKED, self.risk_level, "PPTX path is outside allowed roots.")
        if not path.exists() or path.suffix.lower() != ".pptx":
            return ToolResult(self.name, ActionStatus.FAILED, self.risk_level, "PPTX file does not exist.")
        try:
            presentation = Presentation(str(path))
        except Exception as exc:
            return ToolResult(self.name, ActionStatus.FAILED, self.risk_level, "PPTX file could not be opened.", error=str(exc))
        sample_slides = max(1, min(int(tool_input.get("sample_slides") or 10), 50))
        slides = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if index <= sample_slides:
                slides.append({"index": index, "title": texts[0] if texts else "", "texts": texts[:20]})
        return ToolResult(
            self.name,
            ActionStatus.SUCCEEDED,
            self.risk_level,
            f"Inspected PPTX deck {path}.",
            {"path": str(path), "slide_count": len(presentation.slides), "slides": slides, "source": "pptx_deck_inspect"},
        )


class PresentationPlanCreateTool(Tool):
    def __init__(self) -> None:
        super().__init__(
            name="presentation_plan_create",
            description=(
                "Create a local presentation-design plan with audience, goal, desired action, narrative arc, "
                "slide sequence, visual intent, speaker notes, evidence refs, and review risks."
            ),
            risk_level=RiskLevel.MEDIUM,
            input_schema=object_input_schema(
                {
                    "filename": {"type": "string", "description": "Output markdown filename under data_dir/presentations/plans."},
                    "title": {"type": "string"},
                    "audience": {"type": "string"},
                    "goal": {"type": "string"},
                    "desired_action": {"type": "string"},
                    "status": {"type": "string", "enum": ["draft", "ready_for_review", "final"]},
                    "narrative_arc": {"type": "array", "items": {"type": "string"}},
                    "slide_plan": {"type": "array", "items": {"type": "object"}},
                    "design_notes": {"type": "array", "items": {"type": "string"}},
                    "evidence_refs": {"type": "array", "items": {"type": "string"}},
                    "risks": {"type": "array", "items": {"type": "string"}},
                    "reason": {"type": "string"},
                },
                required=["title", "audience", "goal", "reason"],
            ),
            capability_group="office",
        )

    def execute(self, tool_input: dict[str, Any], config: AgentConfig) -> ToolResult:
        normalized = config.normalized()
        title = " ".join(str(tool_input.get("title") or "").split())
        audience = _bounded_text(tool_input.get("audience"))
        goal = _bounded_text(tool_input.get("goal"))
        reason = str(tool_input.get("reason") or "").strip()
        if not title or not audience or not goal or not reason:
            return ToolResult(self.name, ActionStatus.FAILED, self.risk_level, "Title, audience, goal, and reason are required.")
        filename = _safe_filename(str(tool_input.get("filename") or f"presentation-plan-{uuid4().hex[:8]}.md"), ".md")
        markdown_path = (normalized.data_dir / "presentations" / "plans" / filename).resolve()
        if not _is_within(markdown_path, normalized.allowed_write_roots):
            return ToolResult(self.name, ActionStatus.BLOCKED, self.risk_level, "Presentation plan path is outside allowed write roots.")
        artifact = _presentation_plan_artifact(tool_input, title=title, audience=audience, goal=goal, reason=reason, markdown_path=markdown_path)
        markdown = _render_presentation_plan(artifact)
        if config.dry_run:
            return ToolResult(self.name, ActionStatus.SKIPPED, self.risk_level, f"Dry run: would create presentation plan {markdown_path}.", {"path": str(markdown_path), "artifact": artifact})
        markdown_path.parent.mkdir(parents=True, exist_ok=True)
        markdown_path.write_text(markdown, encoding="utf-8")
        metadata_path = markdown_path.with_suffix(".json")
        metadata_path.write_text(json.dumps(artifact, indent=2, ensure_ascii=False, sort_keys=True) + "\n", encoding="utf-8")
        return ToolResult(
            self.name,
            ActionStatus.SUCCEEDED,
            self.risk_level,
            f"Created presentation plan {markdown_path}.",
            {
                "path": str(markdown_path),
                "metadata_path": str(metadata_path),
                "presentation_plan_id": artifact["presentation_plan_id"],
                "status": artifact["status"],
                "slide_count": len(artifact["slide_plan"]),
                "evidence_ref_count": len(artifact["evidence_refs"]),
                "source": "presentation_plan_create",
            },
        )


class PresentationPlanInspectTool(Tool):
    def __init__(self) -> None:
        super().__init__(
            name="presentation_plan_inspect",
            description="Inspect a local presentation-design plan for audience, goal, slide count, evidence refs, review risks, and preview text.",
            risk_level=RiskLevel.LOW,
            input_schema=object_input_schema(
                {"path": {"type": "string", "description": "Workspace-relative or allowed absolute presentation plan markdown path."}},
                required=["path"],
            ),
            capability_group="office",
        )

    def execute(self, tool_input: dict[str, Any], config: AgentConfig) -> ToolResult:
        normalized = config.normalized()
        path = _resolve_office_path(normalized, str(tool_input.get("path") or ""), subdir="presentations/plans", suffix=".md")
        if not _is_within(path, normalized.allowed_read_roots + normalized.allowed_write_roots):
            return ToolResult(self.name, ActionStatus.BLOCKED, self.risk_level, "Presentation plan path is outside allowed roots.")
        if not path.exists() or path.suffix.lower() != ".md":
            return ToolResult(self.name, ActionStatus.FAILED, self.risk_level, "Presentation plan file does not exist.")
        metadata = _load_json_sidecar(path.with_suffix(".json"))
        text = path.read_text(encoding="utf-8")
        return ToolResult(
            self.name,
            ActionStatus.SUCCEEDED,
            self.risk_level,
            f"Inspected presentation plan {path}.",
            {
                "path": str(path),
                "metadata_path": str(path.with_suffix(".json")) if path.with_suffix(".json").exists() else "",
                "presentation_plan_id": metadata.get("presentation_plan_id", ""),
                "title": metadata.get("title", ""),
                "audience": metadata.get("audience", ""),
                "goal": metadata.get("goal", ""),
                "status": metadata.get("status", ""),
                "slide_count": len(metadata.get("slide_plan", [])) if isinstance(metadata.get("slide_plan"), list) else 0,
                "evidence_ref_count": len(metadata.get("evidence_refs", [])) if isinstance(metadata.get("evidence_refs"), list) else 0,
                "risk_count": len(metadata.get("risks", [])) if isinstance(metadata.get("risks"), list) else 0,
                "preview": text[:4000],
                "source": "presentation_plan_inspect",
            },
        )


def default_office_tools() -> dict[str, Tool]:
    tools: list[Tool] = [
        DocxDocumentCreateTool(),
        DocxDocumentInspectTool(),
        PresentationPlanCreateTool(),
        PresentationPlanInspectTool(),
        PptxDeckCreateTool(),
        PptxDeckInspectTool(),
    ]
    return {tool.name: tool for tool in tools}


def _presentation_plan_artifact(tool_input: dict[str, Any], *, title: str, audience: str, goal: str, reason: str, markdown_path: Path) -> dict[str, Any]:
    status = str(tool_input.get("status") or "draft").strip().lower()
    if status not in {"draft", "ready_for_review", "final"}:
        status = "draft"
    return {
        "presentation_plan_id": f"presentation-plan-{uuid4().hex[:12]}",
        "created_at": datetime.now(timezone.utc).isoformat(),
        "title": title,
        "audience": audience,
        "goal": goal,
        "desired_action": _bounded_text(tool_input.get("desired_action")),
        "status": status,
        "narrative_arc": _string_list(tool_input.get("narrative_arc"))[:OFFICE_MAX_SLIDES],
        "slide_plan": _slide_plan(tool_input.get("slide_plan")),
        "design_notes": _string_list(tool_input.get("design_notes"))[:OFFICE_MAX_SLIDES],
        "evidence_refs": _string_list(tool_input.get("evidence_refs"))[:OFFICE_MAX_SLIDES],
        "risks": _string_list(tool_input.get("risks"))[:OFFICE_MAX_SLIDES],
        "reason": reason,
        "path": str(markdown_path),
        "safety_note": "Presentation plans are local design artifacts. Claims, metrics, logos, and brand rules must be verified before final deck use.",
    }


def _render_presentation_plan(artifact: dict[str, Any]) -> str:
    lines = [
        f"# {artifact['title']}",
        "",
        f"Audience: {artifact['audience']}",
        f"Goal: {artifact['goal']}",
        f"Desired action: {artifact['desired_action']}",
        f"Status: {artifact['status']}",
        "",
    ]
    _append_list(lines, "Narrative Arc", artifact["narrative_arc"])
    if artifact["slide_plan"]:
        lines.extend(["## Slide Plan", "", "| # | Title | Purpose | Key Points | Visual | Speaker Notes | Evidence |", "| --- | --- | --- | --- | --- | --- | --- |"])
        for index, slide in enumerate(artifact["slide_plan"], start=1):
            lines.append(
                f"| {index} | {slide['title']} | {slide['purpose']} | {'; '.join(slide['key_points'])} | {slide['visual']} | {slide['speaker_notes']} | {'; '.join(slide['evidence_refs'])} |"
            )
        lines.append("")
    _append_list(lines, "Design Notes", artifact["design_notes"])
    _append_list(lines, "Evidence References", artifact["evidence_refs"])
    _append_list(lines, "Risks And Review Items", artifact["risks"])
    lines.extend(["## Safety Note", "", artifact["safety_note"], "", f"Created: {artifact['created_at']}"])
    return "\n".join(lines) + "\n"


def _slide_plan(value: Any) -> list[dict[str, Any]]:
    slides = []
    if not isinstance(value, list):
        return slides
    for raw in value[:OFFICE_MAX_SLIDES]:
        if not isinstance(raw, dict):
            continue
        title = _bounded_text(raw.get("title"))
        purpose = _bounded_text(raw.get("purpose"))
        key_points = _string_list(raw.get("key_points") or raw.get("bullets"))[:20]
        body = _bounded_text(raw.get("body"))
        if body and not key_points:
            key_points = [body]
        if title or purpose or key_points:
            slides.append(
                {
                    "title": title or "Untitled Slide",
                    "purpose": purpose,
                    "key_points": key_points,
                    "visual": _bounded_text(raw.get("visual") or raw.get("visual_intent")),
                    "speaker_notes": _bounded_text(raw.get("speaker_notes") or raw.get("notes")),
                    "evidence_refs": _string_list(raw.get("evidence_refs"))[:20],
                }
            )
    return slides


def _write_docx(path: Path, *, title: str, sections: list[Any]) -> None:
    body = [_docx_paragraph(title, style="Title")]
    for raw in sections[:100]:
        if not isinstance(raw, dict):
            continue
        heading = str(raw.get("heading") or "").strip()
        if heading:
            body.append(_docx_paragraph(heading, style="Heading1"))
        for paragraph in _string_list(raw.get("paragraphs"))[:OFFICE_MAX_PARAGRAPHS]:
            body.append(_docx_paragraph(paragraph))
        for bullet in _string_list(raw.get("bullets"))[:OFFICE_MAX_PARAGRAPHS]:
            body.append(_docx_paragraph(bullet, style="ListParagraph"))
        for table in _tables(raw.get("tables")):
            body.append(_docx_table(table))
    document_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        f"<w:body>{''.join(body)}<w:sectPr><w:pgSz w:w=\"12240\" w:h=\"15840\"/><w:pgMar w:top=\"1440\" w:right=\"1440\" w:bottom=\"1440\" w:left=\"1440\"/></w:sectPr></w:body>"
        "</w:document>"
    )
    with ZipFile(path, "w", ZIP_DEFLATED) as package:
        package.writestr("[Content_Types].xml", _docx_content_types())
        package.writestr("_rels/.rels", _docx_rels())
        package.writestr("word/document.xml", document_xml)
        package.writestr("word/styles.xml", _docx_styles())


def _inspect_docx(path: Path, *, sample_paragraphs: int) -> dict[str, Any]:
    with ZipFile(path) as package:
        document_xml = package.read("word/document.xml")
    root = ET.fromstring(document_xml)
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    paragraphs = []
    for paragraph in root.findall(".//w:p", ns):
        text = "".join(node.text or "" for node in paragraph.findall(".//w:t", ns)).strip()
        if text:
            paragraphs.append(text)
    tables = root.findall(".//w:tbl", ns)
    return {
        "paragraph_count": len(paragraphs),
        "table_count": len(tables),
        "sample_paragraphs": paragraphs[:sample_paragraphs],
        "text_preview": "\n".join(paragraphs[:sample_paragraphs])[:2000],
    }


def _docx_paragraph(text: str, *, style: str = "") -> str:
    style_xml = f'<w:pPr><w:pStyle w:val="{escape(style)}"/></w:pPr>' if style else ""
    return f"<w:p>{style_xml}<w:r><w:t>{escape(str(text))}</w:t></w:r></w:p>"


def _docx_table(rows: list[list[str]]) -> str:
    xml_rows = []
    for row in rows[:OFFICE_MAX_TABLE_ROWS]:
        cells = "".join(f"<w:tc><w:p><w:r><w:t>{escape(str(cell))}</w:t></w:r></w:p></w:tc>" for cell in row[:20])
        xml_rows.append(f"<w:tr>{cells}</w:tr>")
    return f"<w:tbl>{''.join(xml_rows)}</w:tbl>"


def _docx_content_types() -> str:
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
        '<Override PartName="/word/styles.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.styles+xml"/>'
        "</Types>"
    )


def _docx_rels() -> str:
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>'
        "</Relationships>"
    )


def _docx_styles() -> str:
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:styles xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        '<w:style w:type="paragraph" w:styleId="Title"><w:name w:val="Title"/></w:style>'
        '<w:style w:type="paragraph" w:styleId="Heading1"><w:name w:val="heading 1"/></w:style>'
        '<w:style w:type="paragraph" w:styleId="ListParagraph"><w:name w:val="List Paragraph"/></w:style>'
        "</w:styles>"
    )


def _add_title_slide(presentation: Any, title: str) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Prepared by Humungousaur"


def _add_content_slide(presentation: Any, spec: dict[str, Any]) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = str(spec.get("title") or "Untitled Slide")
    placeholder = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    bullets = _string_list(spec.get("bullets"))
    body = str(spec.get("body") or "").strip()
    if placeholder is not None:
        text_frame = placeholder.text_frame
        text_frame.clear()
        items = bullets or ([body] if body else [])
        for index, item in enumerate(items[:20]):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = item
            paragraph.level = 0


def _docx_summary(title: str, sections: list[Any]) -> dict[str, Any]:
    paragraph_count = 1
    bullet_count = 0
    table_count = 0
    for section in sections:
        if not isinstance(section, dict):
            continue
        paragraph_count += 1 if str(section.get("heading") or "").strip() else 0
        paragraph_count += len(_string_list(section.get("paragraphs")))
        bullets = _string_list(section.get("bullets"))
        bullet_count += len(bullets)
        paragraph_count += len(bullets)
        table_count += len(_tables(section.get("tables")))
    return {"title": title, "section_count": len(sections), "paragraph_count": paragraph_count, "bullet_count": bullet_count, "table_count": table_count}


def _pptx_summary(title: str, slides: list[Any]) -> dict[str, Any]:
    return {"title": title, "slide_count": len([slide for slide in slides if isinstance(slide, dict)]) + 1}


def _append_list(lines: list[str], title: str, items: list[str]) -> None:
    if not items:
        return
    lines.extend([f"## {title}", ""])
    for item in items:
        lines.append(f"- {item}")
    lines.append("")


def _bounded_text(value: Any) -> str:
    return " ".join(str(value or "").split())[:OFFICE_MAX_TEXT_CHARS]


def _load_json_sidecar(path: Path) -> dict[str, Any]:
    if not path.exists():
        return {}
    try:
        loaded = json.loads(path.read_text(encoding="utf-8"))
    except json.JSONDecodeError:
        return {}
    return loaded if isinstance(loaded, dict) else {}


def _string_list(value: Any) -> list[str]:
    if isinstance(value, str):
        return [value] if value.strip() else []
    if not isinstance(value, list):
        return []
    return [str(item).strip() for item in value if str(item).strip()]


def _tables(value: Any) -> list[list[list[str]]]:
    if not isinstance(value, list):
        return []
    tables: list[list[list[str]]] = []
    for table in value[:20]:
        rows = table.get("rows") if isinstance(table, dict) else table
        if not isinstance(rows, list):
            continue
        clean_rows = []
        for row in rows[:OFFICE_MAX_TABLE_ROWS]:
            if isinstance(row, list):
                clean_rows.append([str(cell) for cell in row[:20]])
        if clean_rows:
            tables.append(clean_rows)
    return tables


def _safe_filename(value: str, suffix: str) -> str:
    name = Path(value).name.strip() or f"artifact{suffix}"
    if not name.lower().endswith(suffix):
        name += suffix
    stem = "".join(char if char.isalnum() or char in ("-", "_", ".") else "-" for char in Path(name).stem).strip(".-")
    return f"{stem or 'artifact'}{suffix}"


def _resolve_office_path(config: AgentConfig, raw_path: str, *, subdir: str, suffix: str) -> Path:
    path = Path(raw_path).expanduser()
    if not path.is_absolute():
        path = config.workspace / path
        if not path.exists():
            data_path = config.data_dir / raw_path
            if data_path.exists():
                path = data_path
            else:
                artifact_path = config.data_dir / subdir / Path(raw_path).name
                if artifact_path.exists():
                    path = artifact_path
    if not path.suffix:
        path = path.with_suffix(suffix)
    return path.resolve()


def _is_within(path: Path, roots: tuple[Path, ...]) -> bool:
    return any(path == root or root in path.parents for root in roots)
